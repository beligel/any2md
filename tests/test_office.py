"""Unit-тесты для office-экстрактора."""

from pathlib import Path

import pytest

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

from any2md.extractors.office import extract_docx, extract_pptx, extract_xlsx
from any2md.extractors.registry import ExtractionContext


@pytest.mark.skipif(not HAS_DOCX, reason="python-docx не установлен")
def test_docx_heading_and_bold(tmp_path):
    doc = docx.Document()
    doc.add_heading("Заголовок 1", level=1)
    p = doc.add_paragraph()
    p.add_run("жирный").bold = True
    p.add_run(" и ")
    p.add_run("курсив").italic = True

    src = tmp_path / "test.docx"
    doc.save(src)

    ctx = ExtractionContext()
    text = extract_docx(src, ctx)
    assert "# Заголовок 1" in text
    assert "**жирный**" in text
    assert "*курсив*" in text


@pytest.mark.skipif(not HAS_DOCX, reason="python-docx не установлен")
def test_docx_empty_paragraph_skipped(tmp_path):
    """AC: пустые параграфы не попадают в вывод."""
    doc = docx.Document()
    doc.add_paragraph("   ")
    doc.add_paragraph("Реальный текст")

    src = tmp_path / "test.docx"
    doc.save(src)

    ctx = ExtractionContext()
    text = extract_docx(src, ctx)
    assert "Реальный текст" in text
    assert text.strip().count("\n\n") == 0  # только один параграф


@pytest.mark.skipif(not HAS_DOCX, reason="python-docx не установлен")
def test_docx_table_markdown(tmp_path):
    """AC: таблицы конвертируются в markdown-таблицы."""
    doc = docx.Document()
    doc.add_paragraph("Текст до таблицы")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Заголовок A"
    table.cell(0, 1).text = "Заголовок B"
    table.cell(1, 0).text = "Ячейка 1"
    table.cell(1, 1).text = "Ячейка 2"

    src = tmp_path / "test.docx"
    doc.save(src)

    ctx = ExtractionContext()
    text = extract_docx(src, ctx)
    assert "| Заголовок A | Заголовок B |" in text
    assert "| Ячейка 1 | Ячейка 2 |" in text


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx не установлен")
def test_pptx_title_and_notes(tmp_path):
    """AC: заголовок слайда и заметки спикера извлекаются.

    Регрессия: в 0.2.1 extract_pptx был определён дважды —
    старое определение без title/notes перекрывало новое.
    """
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Заголовок презентации"

    # Заметки спикера
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Важные заметки"

    src = tmp_path / "test.pptx"
    prs.save(src)

    ctx = ExtractionContext()
    text = extract_pptx(src, ctx)
    assert "Заголовок презентации" in text
    # AC-2: заметки спикера должны извлекаться
    assert "Важные заметки" in text, (
        "Заметки спикера не извлекаются — "
        "вероятно, активна старая (упрощённая) версия extract_pptx"
    )
    # AC-3: заголовок должен быть выделен как **Заголовок:** (новый формат)
    assert "**Заголовок:**" in text, (
        "Заголовок слайда не отформатирован — "
        "активна старая версия extract_pptx без форматирования title"
    )


@pytest.mark.skipif(not HAS_XLSX, reason="openpyxl не установлен")
def test_xlsx_markdown_table(tmp_path):
    """AC: XLSX конвертируется в markdown-таблицу с разделителем."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Данные"
    ws.append(["Имя", "Возраст"])
    ws.append(["Алиса", 30])
    ws.append(["Боб", 25])

    src = tmp_path / "test.xlsx"
    wb.save(src)

    ctx = ExtractionContext()
    text = extract_xlsx(src, ctx)
    assert "## Лист: Данные" in text
    assert "| Имя | Возраст |" in text
    assert "| Алиса | 30 |" in text
    # Разделитель markdown-таблицы
    assert "---" in text


# ---------------------------------------------------------------------------
# DOCX: списки (List Bullet / List Number)
# ---------------------------------------------------------------------------

@pytest.mark.skipif(not HAS_DOCX, reason="python-docx не установлен")
def test_docx_bullet_list(tmp_path):
    """AC: List Bullet стили конвертируются в '- item'."""
    doc = docx.Document()
    for item in ["Яблоко", "Банан", "Вишня"]:
        doc.add_paragraph(item, style="List Bullet")

    src = tmp_path / "test.docx"
    doc.save(src)

    ctx = ExtractionContext()
    text = extract_docx(src, ctx)
    assert "- Яблоко" in text
    assert "- Банан" in text
    assert "- Вишня" in text


@pytest.mark.skipif(not HAS_DOCX, reason="python-docx не установлен")
def test_docx_numbered_list(tmp_path):
    """AC: List Number стили конвертируются в '1. item'."""
    doc = docx.Document()
    for item in ["Первый", "Второй"]:
        doc.add_paragraph(item, style="List Number")

    src = tmp_path / "test.docx"
    doc.save(src)

    ctx = ExtractionContext()
    text = extract_docx(src, ctx)
    assert "1. Первый" in text
    assert "1. Второй" in text


@pytest.mark.skipif(not HAS_DOCX, reason="python-docx не установлен")
def test_docx_underline(tmp_path):
    """AC: underline runs конвертируются в <u>."""
    doc = docx.Document()
    p = doc.add_paragraph()
    run = p.add_run("подчёркнутый")
    run.underline = True

    src = tmp_path / "test.docx"
    doc.save(src)

    ctx = ExtractionContext()
    text = extract_docx(src, ctx)
    assert "<u>подчёркнутый</u>" in text


# ---------------------------------------------------------------------------
# XLSX: edge cases
# ---------------------------------------------------------------------------

@pytest.mark.skipif(not HAS_XLSX, reason="openpyxl не установлен")
def test_xlsx_empty_sheet(tmp_path):
    """AC: пустой лист не добавляет markdown-таблицу."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Пустой"
    # Не добавляем строки

    src = tmp_path / "test.xlsx"
    wb.save(src)

    ctx = ExtractionContext()
    text = extract_xlsx(src, ctx)
    assert "## Лист: Пустой" in text
    # Нет markdown-таблицы (нет строк)
    assert "---" not in text


@pytest.mark.skipif(not HAS_XLSX, reason="openpyxl не установлен")
def test_xlsx_multiple_sheets(tmp_path):
    """AC: несколько листов — каждый с заголовком."""
    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Лист1"
    ws1.append(["A", "B"])

    ws2 = wb.create_sheet("Лист2")
    ws2.append(["X", "Y"])

    src = tmp_path / "test.xlsx"
    wb.save(src)

    ctx = ExtractionContext()
    text = extract_xlsx(src, ctx)
    assert "## Лист: Лист1" in text
    assert "## Лист: Лист2" in text


# ---------------------------------------------------------------------------
# PPTX: edge cases
# ---------------------------------------------------------------------------

@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx не установлен")
def test_pptx_body_without_title(tmp_path):
    """AC: слайд без title placeholder — только body текст."""
    prs = Presentation()
    # Используем layout без title (Blank)
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(100, 100, 500, 200)
    tf = txBox.text_frame
    tf.text = "Только body текст"

    src = tmp_path / "test.pptx"
    prs.save(src)

    ctx = ExtractionContext()
    text = extract_pptx(src, ctx)
    assert "Только body текст" in text
    # Нет заголовка — не должно быть **Заголовок:**
    assert "**Заголовок:**" not in text


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx не установлен")
def test_pptx_no_notes(tmp_path):
    """AC: слайд без заметок — нет секции заметок."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(100, 100, 500, 200)
    tf = txBox.text_frame
    tf.text = "Без заметок"

    src = tmp_path / "test.pptx"
    prs.save(src)

    ctx = ExtractionContext()
    text = extract_pptx(src, ctx)
    assert "Заметки" not in text


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx не установлен")
def test_pptx_empty_notes(tmp_path):
    """AC: пустые заметки спикера — не попадают в вывод."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Тест"
    # Заметки пустые по умолчанию

    src = tmp_path / "test.pptx"
    prs.save(src)

    ctx = ExtractionContext()
    text = extract_pptx(src, ctx)
    assert "Тест" in text
    # Пустые заметки не должны добавлять секцию
    lines = [l for l in text.split("\n") if "Заметки к слайду" in l]
    assert len(lines) == 0


@pytest.mark.skipif(not HAS_PPTX, reason="python-pptx не установлен")
def test_pptx_multiple_slides(tmp_path):
    """AC: несколько слайдов — каждый с заголовком Слайд N."""
    prs = Presentation()
    for i in range(3):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(100, 100, 500, 200)
        txBox.text_frame.text = f"Слайд {i+1}"

    src = tmp_path / "test.pptx"
    prs.save(src)

    ctx = ExtractionContext()
    text = extract_pptx(src, ctx)
    assert "## Слайд 1" in text
    assert "## Слайд 2" in text
    assert "## Слайд 3" in text
