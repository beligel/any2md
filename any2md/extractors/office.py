"""Экстракторы для офисных форматов: DOCX, XLSX, PPTX."""

from pathlib import Path

from .registry import ExtractionContext, register_extractor


_STYLE_TO_MD = {
    "Heading 1": "# ",
    "Heading 2": "## ",
    "Heading 3": "### ",
    "Heading 4": "#### ",
    "Heading 5": "##### ",
    "Heading 6": "###### ",
    "Title": "# ",
    "Subtitle": "## ",
}


def _para_to_md(para) -> str:
    """Конвертирует docx Paragraph в markdown с сохранением базового форматирования."""
    text = para.text
    if not text.strip():
        return ""
    style = para.style.name if para.style else ""
    prefix = _STYLE_TO_MD.get(style, "")
    # Жирный/курсив/код внутри runs
    chunks = []
    for run in para.runs:
        chunk = run.text
        if not chunk:
            continue
        if run.bold:
            chunk = f"**{chunk}**"
        if run.italic:
            chunk = f"*{chunk}*"
        if run.underline:
            chunk = f"<u>{chunk}</u>"
        if run.font and run.font.name and run.font.name.lower() in {"courier new", "consolas", "monospace"}:
            chunk = f"`{chunk}`"
        chunks.append(chunk)
    return prefix + "".join(chunks)


@register_extractor(".docx")
@register_extractor("application/vnd.openxmlformats-officedocument.wordprocessingml.document")
def extract_docx(path: Path, ctx: ExtractionContext) -> str:
    try:
        import docx
    except ImportError:
        raise RuntimeError("Для .docx нужно установить python-docx")
    doc = docx.Document(path)
    parts = []

    prev_list = None
    for para in doc.paragraphs:
        md = _para_to_md(para)
        if not md:
            continue
        # Список
        if para.style and para.style.name.startswith("List"):
            level = 0
            try:
                level = int(para.style.name.split()[-1]) - 1 if para.style.name.split()[-1].isdigit() else 0
            except Exception:
                pass
            indent = "  " * max(level, 0)
            # Маркер/нумератор — упрощённо
            if "Number" in para.style.name:
                md = f"{indent}1. {md}"
            else:
                md = f"{indent}- {md}"
        parts.append(md)

    for table in doc.tables:
        parts.append("")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append("| " + " | ".join(cells) + " |")
    return "\n\n".join(parts)


@register_extractor(".xlsx")
@register_extractor("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
def extract_xlsx(path: Path, ctx: ExtractionContext) -> str:
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("Для .xlsx нужно установить openpyxl")
    wb = openpyxl.load_workbook(path, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"## Лист: {sheet.title}")
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        md_rows = ["| " + " | ".join(str(c or "") for c in row) + " |" for row in rows]
        sep = "|" + "|".join([" --- " for _ in rows[0]]) + "|"
        parts.append("\n".join([md_rows[0], sep] + md_rows[1:]))
    return "\n\n".join(parts)


@register_extractor(".pptx")
@register_extractor("application/vnd.openxmlformats-officedocument.presentationml.presentation")
def extract_pptx(path: Path, ctx: ExtractionContext) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("Для .pptx нужно установить python-pptx")
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Слайд {i}")
        title = None
        body = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape.is_placeholder and getattr(shape, "placeholder_format", None):
                    pf = shape.placeholder_format
                    if pf.type in {1, 3}:  # TITLE / CENTER_TITLE
                        title = shape.text.strip()
                        continue
                body.append(shape.text.strip())
        if title:
            parts.append(f"**Заголовок:** {title}")
        parts.extend(body)
        # Заметки спикера
        if slide.has_notes_slide and slide.notes_slide:
            notes_text_frame = slide.notes_slide.notes_text_frame
            if notes_text_frame and notes_text_frame.text.strip():
                parts.append(f"### Заметки к слайду {i}\n{notes_text_frame.text.strip()}")
    return "\n\n".join(parts)


@register_extractor(".xlsx")
@register_extractor("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
def extract_xlsx(path: Path, ctx: ExtractionContext) -> str:
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("Для .xlsx нужно установить openpyxl")
    wb = openpyxl.load_workbook(path, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"## Лист: {sheet.title}")
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        md_rows = ["| " + " | ".join(str(c or "") for c in row) + " |" for row in rows]
        sep = "|" + "|".join([" --- " for _ in rows[0]]) + "|"
        parts.append("\n".join([md_rows[0], sep] + md_rows[1:]))
    return "\n\n".join(parts)


@register_extractor(".pptx")
@register_extractor("application/vnd.openxmlformats-officedocument.presentationml.presentation")
def extract_pptx(path: Path, ctx: ExtractionContext) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("Для .pptx нужно установить python-pptx")
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Слайд {i}")
        title = None
        body = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if shape.is_placeholder and getattr(shape, "placeholder_format", None):
                    pf = shape.placeholder_format
                    if pf.type in {1, 3}:  # TITLE / CENTER_TITLE
                        title = shape.text.strip()
                        continue
                body.append(shape.text.strip())
        if title:
            parts.append(f"**Заголовок:** {title}")
        parts.extend(body)
        # Заметки спикера
        if slide.has_notes_slide and slide.notes_slide:
            notes_text_frame = slide.notes_slide.notes_text_frame
            if notes_text_frame and notes_text_frame.text.strip():
                parts.append(f"### Заметки к слайду {i}\n{notes_text_frame.text.strip()}")
    return "\n\n".join(parts)
